"""
Lightweight text extractors for common document formats.

This module provides fast, dependency-light text extraction without requiring
the unstructured library. Uses direct format-specific libraries:
- openpyxl for Excel (.xlsx, .xls)
- python-docx for Word (.docx)
- python-pptx for PowerPoint (.pptx)
- PyMuPDF for PDF (.pdf)
- BeautifulSoup/lxml for HTML (.html, .htm)
- chardet for text files (.txt, .json, .csv, .md, etc.)

Falls back to unstructured only for rare formats like .eml, .epub.
"""

import io
import json as json_lib
import logging
import re
from typing import Tuple

logger = logging.getLogger(__name__)


def extract_xlsx_text(file_bytes: bytes, filename: str = "") -> str:
    """
    Extract text from Excel files using openpyxl.

    Based on approach from Danswer/Onyx - simpler and faster than unstructured.
    """
    try:
        import openpyxl
    except ImportError:
        raise ImportError("openpyxl is required for XLSX parsing. Install with: pip install openpyxl")

    try:
        workbook = openpyxl.load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
    except Exception as e:
        logger.warning(f"Failed to load workbook {filename}: {e}")
        raise

    text_content = []

    for sheet in workbook.worksheets:
        rows = []
        num_empty_consecutive_rows = 0

        for row in sheet.iter_rows(min_row=1, values_only=True):
            row_str = ",".join(str(cell or "") for cell in row)

            # Only add the row if there are any values in the cells
            if len(row_str) >= len(row):
                rows.append(row_str)
                num_empty_consecutive_rows = 0
            else:
                num_empty_consecutive_rows += 1

            # Handle massive excel sheets with mostly empty cells
            if num_empty_consecutive_rows > 100:
                logger.debug(f"Found {num_empty_consecutive_rows} empty rows in {filename}, " "skipping rest of sheet")
                break

        sheet_str = "\n".join(rows)
        if sheet_str.strip():  # Only add non-empty sheets
            text_content.append(f"Sheet: {sheet.title}\n{sheet_str}")

    workbook.close()
    return "\n\n".join(text_content)


def extract_docx_text(file_bytes: bytes, filename: str = "") -> str:
    """
    Extract text from Word documents using python-docx.

    Extracts paragraphs and table content.
    """
    try:
        import docx
    except ImportError:
        raise ImportError("python-docx is required for DOCX parsing. Install with: pip install python-docx")

    try:
        doc = docx.Document(io.BytesIO(file_bytes))
    except Exception as e:
        logger.warning(f"Failed to load DOCX {filename}: {e}")
        raise

    paragraphs = []

    # Extract text from paragraphs and tables
    for element in doc.element.body:
        # Handle paragraphs
        if element.tag.endswith("p"):
            para = next((p for p in doc.paragraphs if p._element == element), None)
            if para and para.text.strip():
                paragraphs.append(para.text.strip())

        # Handle tables
        elif element.tag.endswith("tbl"):
            table = next((t for t in doc.tables if t._element == element), None)
            if table:
                table_text = []
                for row in table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells)
                    if row_text.strip():
                        table_text.append(row_text)
                if table_text:
                    paragraphs.append("\n".join(table_text))

    return "\n\n".join(paragraphs)


def extract_pptx_text(file_bytes: bytes, filename: str = "") -> str:
    """
    Extract text from PowerPoint presentations using python-pptx.

    Extracts text from all slides and shapes.
    """
    try:
        import pptx
    except ImportError:
        raise ImportError("python-pptx is required for PPTX parsing. Install with: pip install python-pptx")

    try:
        presentation = pptx.Presentation(io.BytesIO(file_bytes))
    except Exception as e:
        logger.warning(f"Failed to load PPTX {filename}: {e}")
        raise

    text_content = []

    for slide_num, slide in enumerate(presentation.slides, start=1):
        slide_text = [f"\n=== Slide {slide_num} ===\n"]

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text.strip())

        if len(slide_text) > 1:  # Has more than just the header
            text_content.append("\n".join(slide_text))

    return "\n\n".join(text_content)


def detect_encoding(file_bytes: bytes) -> str:
    """Detect encoding of text file using chardet."""
    try:
        import chardet
    except ImportError:
        logger.warning("chardet not available, defaulting to utf-8")
        return "utf-8"

    # Check first 50KB for encoding detection
    sample = file_bytes[:50000]
    detected = chardet.detect(sample)
    encoding = detected.get("encoding") or "utf-8"
    logger.debug(f"Detected encoding: {encoding} (confidence: {detected.get('confidence', 0):.2f})")
    return encoding


def extract_text_file(file_bytes: bytes, filename: str = "") -> str:
    """
    Extract text from plain text files (.txt, .md, .csv, etc.).

    Detects encoding automatically and handles various text formats.
    """
    encoding = detect_encoding(file_bytes)

    try:
        text = file_bytes.decode(encoding)
        return text
    except UnicodeDecodeError:
        # Fallback to utf-8 with error handling
        logger.warning(f"Failed to decode {filename} with {encoding}, trying utf-8 with replace")
        try:
            text = file_bytes.decode("utf-8", errors="replace")
            return text
        except Exception as e:
            logger.error(f"Failed to decode {filename}: {e}")
            raise


def extract_json_file(file_bytes: bytes, filename: str = "") -> str:
    """
    Extract text from JSON files.

    Pretty-prints JSON for better readability.
    """
    encoding = detect_encoding(file_bytes)

    try:
        # Try to parse as JSON and pretty-print
        json_data = json_lib.loads(file_bytes.decode(encoding))
        return json_lib.dumps(json_data, indent=2, ensure_ascii=False)
    except json_lib.JSONDecodeError:
        # If not valid JSON, just return as text
        logger.warning(f"{filename} is not valid JSON, treating as plain text")
        return extract_text_file(file_bytes, filename)
    except Exception as e:
        logger.error(f"Failed to parse JSON {filename}: {e}")
        raise


def extract_pdf_text(file_bytes: bytes, filename: str = "") -> str:
    """
    Extract text from PDF files using PyMuPDF (fitz).

    PyMuPDF is already a dependency in Morphik (pymupdf==1.25.5).
    """
    try:
        import fitz  # PyMuPDF
    except ImportError:
        raise ImportError("PyMuPDF is required for PDF parsing. Install with: pip install pymupdf")

    try:
        pdf_document = fitz.open("pdf", file_bytes)
    except Exception as e:
        logger.warning(f"Failed to open PDF {filename}: {e}")
        raise

    text_content = []

    for page_num in range(len(pdf_document)):
        page = pdf_document[page_num]
        try:
            page_text = page.get_text()
            if page_text.strip():
                text_content.append(f"=== Page {page_num + 1} ===\n{page_text}")
        except Exception as e:
            logger.warning(f"Failed to extract text from page {page_num + 1} of {filename}: {e}")
            continue

    pdf_document.close()

    if not text_content:
        logger.warning(f"No text extracted from PDF {filename}")
        return ""

    return "\n\n".join(text_content)


def extract_html_text(file_bytes: bytes, filename: str = "") -> str:
    """
    Extract text from HTML files.

    Uses lxml (already a dependency) for parsing, with BeautifulSoup if available.
    """
    encoding = detect_encoding(file_bytes)
    html_content = file_bytes.decode(encoding, errors="replace")

    try:
        from bs4 import BeautifulSoup

        soup = BeautifulSoup(html_content, "lxml")

        # Remove script and style elements
        for script in soup(["script", "style", "meta", "link"]):
            script.decompose()

        # Get text
        text = soup.get_text(separator="\n", strip=True)

        # Clean up whitespace
        lines = (line.strip() for line in text.splitlines())
        chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
        text = "\n".join(chunk for chunk in chunks if chunk)

        return text

    except ImportError:
        logger.warning("BeautifulSoup not available, using simple HTML parser")
        # Fallback: simple regex-based HTML tag removal
        text = re.sub(r"<script[^>]*>.*?</script>", "", html_content, flags=re.DOTALL | re.IGNORECASE)
        text = re.sub(r"<style[^>]*>.*?</style>", "", text, flags=re.DOTALL | re.IGNORECASE)
        text = re.sub(r"<[^>]+>", " ", text)
        text = re.sub(r"\s+", " ", text)
        return text.strip()

    except Exception as e:
        logger.error(f"Failed to parse HTML {filename}: {e}")
        raise


def extract_text_lightweight(file_bytes: bytes, filename: str) -> Tuple[str, bool]:
    """
    Try to extract text using lightweight parsers based on file extension.

    Returns:
        Tuple of (extracted_text, success)
        - extracted_text: The extracted text (empty string if failed)
        - success: True if extraction succeeded, False if should fallback to unstructured
    """
    extension = filename.lower().rsplit(".", 1)[-1] if "." in filename else ""

    # Plain text extensions that don't need special processing
    text_extensions = {"txt", "md", "mdx", "conf", "log", "csv", "tsv", "xml", "yml", "yaml", "sql"}

    try:
        # PDF files
        if extension == "pdf":
            logger.info(f"Using lightweight PDF parser for {filename}")
            text = extract_pdf_text(file_bytes, filename)
            return text, True

        # HTML files
        elif extension in ("html", "htm"):
            logger.info(f"Using lightweight HTML parser for {filename}")
            text = extract_html_text(file_bytes, filename)
            return text, True

        # Office formats
        elif extension in ("xlsx", "xls"):
            logger.info(f"Using lightweight XLSX parser for {filename}")
            text = extract_xlsx_text(file_bytes, filename)
            return text, True

        elif extension == "docx":
            logger.info(f"Using lightweight DOCX parser for {filename}")
            text = extract_docx_text(file_bytes, filename)
            return text, True

        elif extension == "pptx":
            logger.info(f"Using lightweight PPTX parser for {filename}")
            text = extract_pptx_text(file_bytes, filename)
            return text, True

        # JSON files
        elif extension == "json":
            logger.info(f"Using lightweight JSON parser for {filename}")
            text = extract_json_file(file_bytes, filename)
            return text, True

        # Plain text files
        elif extension in text_extensions:
            logger.info(f"Using lightweight text parser for {filename}")
            text = extract_text_file(file_bytes, filename)
            return text, True

        else:
            # Not a supported lightweight format (only .eml, .epub left)
            return "", False

    except Exception as e:
        logger.warning(
            f"Lightweight parsing failed for {filename}: {e}. " "Will fallback to unstructured if available."
        )
        return "", False
